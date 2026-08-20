"""將 `assembly.assemble_document()` 的組裝結果，渲染成一份教育訓練簡報
PPTX 檔案（Issue #10：投影片序列結構，一張投影片＝一個標題＋條列項目，不
做「整段文字自動轉條列」的處理——`bullets` 已經是後台以條列 UI 輸入的結構
化陣列，這裡只負責照原樣渲染成投影片上的條列文字，不做任何文字分行/斷句
邏輯）。

刻意獨立成一個模組（同 `docx_builder.py` 的既有慣例，而不是塞進
views.py）：渲染邏輯本身不碰 HTTP，純粹是「(assembled dict) -> bytes」的
轉換函式。PPTX 有自己的一份 builder（而不是重用/修改 `docx_builder.py`）
——兩種輸出格式的版面邏輯完全不同（章節樹 vs 投影片序列），沒有值得共用
的排版程式碼，硬要共用只會讓兩邊都變得更難懂（見 Issue #10 Explicit scope
boundaries：「PPTX gets its own renderer」）。

設計判斷（PPTX 專屬，跟 DOCX 不同，誠實標註——Issue #10 沒有明講版面細節）:

1. **不額外加封面投影片**：Issue #1 的 DOCX/PPTX Output Spec 只對 DOCX
   明講需要保留封面／頁首頁尾／頁碼；PPTX 段落只要求「投影片序列，一張
   投影片＝一個標題＋條列項目」。這裡刻意不多加一張非「章節投影片」的
   投影片、不加頁碼/頁尾，讓「投影片數 = 已包含章節數」這個關係維持單純
   可預期，貼近 Issue #10 acceptance criteria 原文「投影片依照後台設定的
   條列內容組成」的字面意思。
2. **版面尺寸採 16:9 寬螢幕**（13.333 x 7.5 吋），而非 python-pptx 預設
   的 4:3——現行簡報軟體/投影設備幾乎都是 16:9，純排版選擇，不影響任何
   「可觀察內容」的驗收判準。
3. **截圖版面**：條列項目固定放在投影片左側（無截圖時佔滿全寬）；若該
   章節有可用的截圖，依序垂直堆疊在投影片右側欄位，圖片下方緊接圖說文
   字——跟 DOCX「內容之後接著圖片＋圖說」的閱讀順序精神一致，但改成左
   右並排（而非上下堆疊），因為投影片是橫向版面，左右並排能讓條列文字
   與圖片同時可見、不需要在同一張投影片內再往下捲動，比照搬 DOCX 的上下
   排列到投影片上更符合簡報排版習慣。
4. **條列符號**：python-pptx 建立的一般文字方塊（非版型內建的
   content placeholder）不會自動套用項目符號樣式，這裡以純文字前綴
   「• 」表示條列——不引入額外的 OXML 項目符號設定，保持渲染邏輯單純，
   同時仍讓輸出檔案「看起來」是條列（而不是一整段連續散文，呼應 Issue
   #10「不做整段文字自動轉條列」的精神：這裡渲染的原本就已經是陣列，只
   是加上視覺上的條列符號）。
"""

import io

from pptx import Presentation
from pptx.util import Inches, Pt

from screenshots.models import ProjectScreenshot

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)

_MARGIN = Inches(0.5)
_TITLE_TOP = Inches(0.3)
_TITLE_HEIGHT = Inches(1.0)

_CONTENT_TOP = _MARGIN + _TITLE_HEIGHT

_BULLETS_ONLY_WIDTH = _SLIDE_WIDTH - (2 * _MARGIN)
_BULLETS_WITH_IMAGE_WIDTH = Inches(7.2)

_IMAGE_COLUMN_LEFT = Inches(8.1)
_IMAGE_COLUMN_WIDTH = Inches(4.7)
_IMAGE_HEIGHT = Inches(2.6)
_IMAGE_GAP = Inches(0.25)
_CAPTION_HEIGHT = Inches(0.45)

_BLANK_LAYOUT_INDEX = 6


def _add_title(slide, title_text: str) -> None:
    box = slide.shapes.add_textbox(
        _MARGIN, _TITLE_TOP, _SLIDE_WIDTH - (2 * _MARGIN), _TITLE_HEIGHT
    )
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.text = title_text
    run = text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(32)
    run.font.bold = True


def _add_bullets(slide, bullets: list, *, has_image: bool) -> None:
    width = _BULLETS_WITH_IMAGE_WIDTH if has_image else _BULLETS_ONLY_WIDTH
    height = _SLIDE_HEIGHT - _CONTENT_TOP - _MARGIN
    box = slide.shapes.add_textbox(_MARGIN, _CONTENT_TOP, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True

    if not bullets:
        return

    text_frame.text = f"• {bullets[0]}"
    text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    for bullet_text in bullets[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = f"• {bullet_text}"
        paragraph.runs[0].font.size = Pt(20)


def _load_image_bytes(shot: "ProjectScreenshot") -> bytes:
    shot.image.open("rb")
    try:
        return shot.image.read()
    finally:
        shot.image.close()


def _add_screenshots(slide, screenshot_entries: list) -> None:
    top = _CONTENT_TOP
    for entry in screenshot_entries:
        if not entry.get("applicable"):
            continue
        screenshot_ref = entry.get("screenshot")
        if not screenshot_ref:
            continue

        shot = ProjectScreenshot.objects.filter(pk=screenshot_ref["id"]).first()
        if shot is None or not shot.image:
            continue

        image_bytes = _load_image_bytes(shot)
        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            _IMAGE_COLUMN_LEFT,
            top,
            width=_IMAGE_COLUMN_WIDTH,
            height=_IMAGE_HEIGHT,
        )
        top += _IMAGE_HEIGHT

        caption_text = entry.get("caption") or shot.original_filename
        caption_box = slide.shapes.add_textbox(
            _IMAGE_COLUMN_LEFT, top, _IMAGE_COLUMN_WIDTH, _CAPTION_HEIGHT
        )
        caption_frame = caption_box.text_frame
        caption_frame.word_wrap = True
        caption_frame.text = f"圖：{caption_text}"
        caption_run = caption_frame.paragraphs[0].runs[0]
        caption_run.font.size = Pt(12)
        caption_run.font.italic = True
        top += _CAPTION_HEIGHT + _IMAGE_GAP


def _has_applicable_screenshot(screenshot_entries: list) -> bool:
    return any(
        entry.get("applicable") and entry.get("screenshot")
        for entry in screenshot_entries
    )


def _add_slide(presentation: Presentation, section: dict) -> None:
    layout = presentation.slide_layouts[_BLANK_LAYOUT_INDEX]
    slide = presentation.slides.add_slide(layout)

    screenshot_entries = section.get("screenshots", [])
    has_image = _has_applicable_screenshot(screenshot_entries)

    _add_title(slide, section["title"])
    _add_bullets(slide, section.get("bullets", []), has_image=has_image)
    if has_image:
        _add_screenshots(slide, screenshot_entries)


def build_pptx(assembled: dict) -> bytes:
    """把 `assemble_document()` 的回傳值（`assembled`）渲染成一份教育訓練
    簡報 PPTX 檔案的完整 bytes。只渲染 `included=True` 的章節，一個章節對
    應一張投影片（同 `docx_builder.build_docx()` 的「只渲染 included 章節」
    既有慣例）。
    """
    presentation = Presentation()
    presentation.slide_width = _SLIDE_WIDTH
    presentation.slide_height = _SLIDE_HEIGHT

    for section in assembled["sections"]:
        if not section["included"]:
            continue
        _add_slide(presentation, section)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
