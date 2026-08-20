"""HTTP-layer tests for T9: PPTX 文件產生（教育訓練簡報，Issue #10）.

同 T8 的 `test_document_generation_api.py`（Issue #9）的測試哲學：HTTP API
層是唯一的測試邊界。對這個二進位輸出格式，HTTP-layer-appropriate 的測試方
式是——呼叫產生端點、拿到回傳的檔案 bytes、用 `python-pptx` 重新解析，斷
言可觀察的結構（投影片數、標題文字、條列文字存在、被排除節點的內容不存
在、有嵌入真正的圖片資料）——這是在驗證 HTTP 端點的真實輸出，不是內部實
作細節。
"""

import io
import shutil
import struct
import tempfile
import zlib

from django.contrib.auth.models import User
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import Client, TestCase, override_settings
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from features.models import FeatureNode, FeatureNodeContent
from projects.models import Project
from screenshots.models import ProjectScreenshot, ScreenshotRequirement
from selections.models import ProjectFeatureExclusion, ProjectFeatureSelection

from ..models import GeneratedDocument

_TEST_MEDIA_ROOT = tempfile.mkdtemp(prefix="gen_doc_documents_pptx_test_")


def tearDownModule():
    shutil.rmtree(_TEST_MEDIA_ROOT, ignore_errors=True)


# 手刻一張結構完整（IHDR/IDAT/IEND、正確 CRC）的最小合法 PNG——同 T8 既有
# 測試檔案的既有慣例（`python-pptx` 的 `add_picture()` 同樣需要能解析圖片
# 檔頭判斷格式與原生尺寸，不能只是「PNG signature + 隨便填充 bytes」）。
def _png_chunk(chunk_type: bytes, data: bytes) -> bytes:
    return (
        struct.pack(">I", len(data))
        + chunk_type
        + data
        + struct.pack(">I", zlib.crc32(chunk_type + data) & 0xFFFFFFFF)
    )


def _valid_png_bytes(width=20, height=20) -> bytes:
    signature = b"\x89PNG\r\n\x1a\n"
    ihdr = _png_chunk(
        b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    )
    raw_scanlines = b"".join(
        b"\x00" + bytes([255, 0, 0]) * width for _ in range(height)
    )
    idat = _png_chunk(b"IDAT", zlib.compress(raw_scanlines))
    iend = _png_chunk(b"IEND", b"")
    return signature + ihdr + idat + iend


def _valid_png_upload(name="shot.png"):
    return SimpleUploadedFile(name, _valid_png_bytes(), content_type="image/png")


def _generate_url(project_id, document_type="training_deck"):
    return f"/api/projects/{project_id}/documents/{document_type}/generate/"


def _history_url(project_id, document_type="training_deck"):
    return f"/api/projects/{project_id}/documents/{document_type}/history/"


def _status_url():
    return "/api/projects/documents-status/"


def _slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return "\n".join(texts)


def _slide_picture_blob_sizes(slide) -> list:
    sizes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            sizes.append(len(shape.image.blob))
    return sizes


@override_settings(MEDIA_ROOT=_TEST_MEDIA_ROOT)
class PptxContentTests(TestCase):
    """Issue #10 acceptance criterion 1：使用者可將教育訓練簡報輸出成
    PPTX 並下載，投影片依照後台設定的條列內容組成（一張投影片＝一個標題＋
    條列項目）；未勾選/被排除節點的內容不出現。"""

    def setUp(self):
        self.user = User.objects.create_user(username="pm_pptx", password="pw12345")
        self.client.login(username="pm_pptx", password="pw12345")
        self.project = Project.objects.create(
            owner=self.user,
            customer_name="測試客戶",
            project_name="測試專案",
            system_url="https://portal.example.com",
        )

        self.included_root = FeatureNode.objects.create(name="帳號管理")
        self.included_child = FeatureNode.objects.create(
            name="編輯帳號", parent=self.included_root
        )
        self.excluded_child = FeatureNode.objects.create(
            name="刪除帳號確認", parent=self.included_root
        )
        self.unchecked_root = FeatureNode.objects.create(name="報表")

        FeatureNodeContent.objects.create(
            node=self.included_root,
            document_type="training_deck",
            bullets=["可以檢視帳號清單", "系統網址：{{project_variables.system_url}}"],
        )
        FeatureNodeContent.objects.create(
            node=self.included_child,
            document_type="training_deck",
            bullets=["可以編輯帳號基本資料"],
        )
        FeatureNodeContent.objects.create(
            node=self.excluded_child,
            document_type="training_deck",
            bullets=["刪除帳號確認條列內容－不應出現"],
        )
        FeatureNodeContent.objects.create(
            node=self.unchecked_root,
            document_type="training_deck",
            bullets=["報表條列內容－不應出現"],
        )

        ProjectFeatureSelection.objects.create(project=self.project, node=self.included_root)
        ProjectFeatureExclusion.objects.create(project=self.project, node=self.excluded_child)

        self.requirement = ScreenshotRequirement.objects.create(
            node=self.included_root, name="帳號列表"
        )
        ProjectScreenshot.objects.create(
            project=self.project,
            requirement=self.requirement,
            image=_valid_png_upload(),
            original_filename="shot.png",
            document_types=["training_deck"],
            is_custom=False,
            caption="帳號列表畫面截圖",
        )

    def test_training_deck_generation_produces_pptx_with_expected_slides(self):
        response = self.client.post(_generate_url(self.project.id))
        self.assertEqual(response.status_code, 200)
        self.assertEqual(
            response["Content-Type"],
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.assertIn("attachment", response["Content-Disposition"])
        self.assertIn(".pptx", response["Content-Disposition"])
        self.assertTrue(len(response.content) > 0)

        record_id = int(response["X-Generated-Document-Id"])
        record = GeneratedDocument.objects.get(pk=record_id)
        self.assertEqual(record.project_id, self.project.id)
        self.assertEqual(record.document_type, "training_deck")
        self.assertEqual(record.output_format, "pptx")

        presentation = PptxPresentation(io.BytesIO(response.content))
        slides = list(presentation.slides)

        # 兩個已包含的章節（帳號管理、編輯帳號）各自一張投影片——一張投影
        # 片＝一個標題＋條列項目，被排除的子孫節點與未勾選節點都不產生
        # 投影片。
        self.assertEqual(len(slides), 2)

        all_text = "\n".join(_slide_text(slide) for slide in slides)

        # 標題與條列文字都出現，且變數已被替換。
        self.assertIn("帳號管理", all_text)
        self.assertIn("可以檢視帳號清單", all_text)
        self.assertIn("https://portal.example.com", all_text)
        self.assertNotIn("{{project_variables", all_text)
        self.assertIn("編輯帳號", all_text)
        self.assertIn("可以編輯帳號基本資料", all_text)

        # 被排除的子孫節點、未勾選節點的條列內容都不出現。
        self.assertNotIn("刪除帳號確認條列內容", all_text)
        self.assertNotIn("報表條列內容", all_text)

        # 條列符號存在（一張投影片＝一個標題＋條列項目，而非整段散文）。
        self.assertIn("• 可以檢視帳號清單", all_text)

        # 截圖確實被嵌入至少一張投影片：至少一張圖，且是真的圖片資料
        # （非空/非佔位），圖說文字也出現。
        blob_sizes = []
        for slide in slides:
            blob_sizes.extend(_slide_picture_blob_sizes(slide))
        self.assertEqual(len(blob_sizes), 1)
        self.assertGreater(blob_sizes[0], 50)
        self.assertIn("帳號列表畫面截圖", all_text)

    def test_training_deck_content_snapshot_uses_bullets_not_body(self):
        response = self.client.post(_generate_url(self.project.id))
        self.assertEqual(response.status_code, 200)
        record_id = int(response["X-Generated-Document-Id"])
        record = GeneratedDocument.objects.get(pk=record_id)

        node_section = next(
            s
            for s in record.content_snapshot["sections"]
            if s["node"] == self.included_root.id
        )
        self.assertEqual(
            node_section["bullets"],
            ["可以檢視帳號清單", "系統網址：https://portal.example.com"],
        )


@override_settings(MEDIA_ROOT=_TEST_MEDIA_ROOT)
class PptxGenerationBlockedByReadinessTests(TestCase):
    """Issue #9/#10：產生前必須呼叫 `validate_generation_readiness()`，
    缺少必要截圖或變數時擋下產生（PPTX 沿用跟 DOCX 完全同一套 gate，不另
    開一套驗證邏輯）。"""

    def setUp(self):
        self.user = User.objects.create_user(username="pm_pptx_block", password="pw12345")
        self.client.login(username="pm_pptx_block", password="pw12345")
        self.project = Project.objects.create(
            owner=self.user, customer_name="客戶", project_name="專案"
        )
        self.node = FeatureNode.objects.create(name="帳號管理")
        FeatureNodeContent.objects.create(
            node=self.node,
            document_type="training_deck",
            bullets=["系統版本：{{project_variables.system_version}}"],
        )
        self.requirement = ScreenshotRequirement.objects.create(
            node=self.node, name="帳號列表"
        )
        ProjectFeatureSelection.objects.create(project=self.project, node=self.node)

    def test_missing_screenshot_and_variable_blocks_pptx_generation(self):
        response = self.client.post(_generate_url(self.project.id))
        self.assertEqual(response.status_code, 400)
        body = response.json()
        self.assertEqual(len(body["missing_screenshots"]), 1)
        self.assertEqual(body["missing_screenshots"][0]["requirement_id"], self.requirement.id)
        self.assertIn("system_version", body["missing_variables"])
        self.assertEqual(GeneratedDocument.objects.count(), 0)


@override_settings(MEDIA_ROOT=_TEST_MEDIA_ROOT)
class PptxRegenerationHistoryTests(TestCase):
    """Issue #10 acceptance criterion 2：重新產生教育訓練簡報時建立新的歷
    史紀錄，不覆蓋先前紀錄，行為與 DOCX 文件的歷史機制一致——同 T8
    `RegenerationHistoryTests` 的驗證方式，改用 training_deck/PPTX。"""

    def setUp(self):
        self.user = User.objects.create_user(username="pm_pptx_hist", password="pw12345")
        self.client.login(username="pm_pptx_hist", password="pw12345")
        self.project = Project.objects.create(
            owner=self.user, customer_name="舊客戶名稱", project_name="專案"
        )
        self.node = FeatureNode.objects.create(name="登入")
        FeatureNodeContent.objects.create(
            node=self.node, document_type="training_deck", bullets=["登入條列內容"]
        )
        ProjectFeatureSelection.objects.create(project=self.project, node=self.node)

    def _generate(self):
        response = self.client.post(_generate_url(self.project.id))
        self.assertEqual(response.status_code, 200)
        return response

    def test_regenerating_after_customer_name_change_creates_new_row_old_row_unchanged(self):
        first_response = self._generate()
        first_id = int(first_response["X-Generated-Document-Id"])

        self.project.customer_name = "新客戶名稱"
        self.project.save(update_fields=["customer_name"])

        second_response = self._generate()
        second_id = int(second_response["X-Generated-Document-Id"])

        self.assertNotEqual(first_id, second_id)
        self.assertEqual(
            GeneratedDocument.objects.filter(
                project=self.project, document_type="training_deck"
            ).count(),
            2,
        )

        first_record = GeneratedDocument.objects.get(pk=first_id)
        second_record = GeneratedDocument.objects.get(pk=second_id)

        self.assertEqual(
            first_record.content_snapshot["sections"][0]["bullets"], ["登入條列內容"]
        )
        self.assertEqual(
            second_record.content_snapshot["sections"][0]["bullets"], ["登入條列內容"]
        )

        # 先前紀錄實際落地儲存的檔案本身未被覆蓋（重新解析 PPTX bytes 驗
        # 證，而不是只看 DB row 數）——這份簡報渲染不含客戶名稱文字，所以
        # 用「產生時間仍是各自那次的 record」佐證兩筆紀錄各自獨立即可。
        first_record.file.open("rb")
        try:
            stored_old_bytes = first_record.file.read()
        finally:
            first_record.file.close()
        old_presentation = PptxPresentation(io.BytesIO(stored_old_bytes))
        self.assertEqual(len(list(old_presentation.slides)), 1)

        second_record.file.open("rb")
        try:
            stored_new_bytes = second_record.file.read()
        finally:
            second_record.file.close()
        new_presentation = PptxPresentation(io.BytesIO(stored_new_bytes))
        self.assertEqual(len(list(new_presentation.slides)), 1)

        self.assertNotEqual(stored_old_bytes, b"")
        self.assertNotEqual(first_record.file.name, second_record.file.name)

    def test_history_endpoint_returns_all_generations_latest_default(self):
        first = self._generate()
        first_id = int(first["X-Generated-Document-Id"])
        second = self._generate()
        second_id = int(second["X-Generated-Document-Id"])

        response = self.client.get(_history_url(self.project.id))
        self.assertEqual(response.status_code, 200)
        results = response.json()["results"]

        self.assertEqual(len(results), 2)
        self.assertEqual({r["id"] for r in results}, {first_id, second_id})
        self.assertEqual(results[0]["id"], second_id)
        self.assertTrue(results[0]["is_latest"])
        self.assertFalse(results[1]["is_latest"])
        for r in results:
            self.assertEqual(r["output_format"], "pptx")
            self.assertIn("sections", r["content_snapshot"])


@override_settings(MEDIA_ROOT=_TEST_MEDIA_ROOT)
class ProjectListStatusIncludesTrainingDeckTests(TestCase):
    """Issue #10 acceptance criterion 3：專案列表可看到教育訓練簡報目前是
    否已產生，與其他三類文件並列呈現。"""

    def setUp(self):
        self.user = User.objects.create_user(username="pm_pptx_status", password="pw12345")
        self.client.login(username="pm_pptx_status", password="pw12345")
        self.project_a = Project.objects.create(
            owner=self.user, customer_name="A客戶", project_name="A專案"
        )
        self.project_b = Project.objects.create(
            owner=self.user, customer_name="B客戶", project_name="B專案"
        )
        self.node = FeatureNode.objects.create(name="登入")
        FeatureNodeContent.objects.create(
            node=self.node, document_type="training_deck", bullets=["登入條列內容"]
        )
        FeatureNodeContent.objects.create(
            node=self.node, document_type="system_test", body="內容"
        )
        ProjectFeatureSelection.objects.create(project=self.project_a, node=self.node)
        ProjectFeatureSelection.objects.create(project=self.project_b, node=self.node)

    def test_status_includes_training_deck_alongside_other_three_types(self):
        response = self.client.get(_status_url())
        results_by_project = {r["project"]: r for r in response.json()["results"]}

        # 尚未產生任何文件前，training_deck 跟其他三類文件一樣是 False，
        # 且四個欄位並列在同一筆回應中。
        entry_a = results_by_project[self.project_a.id]
        for document_type in ("system_design", "system_test", "system_install", "training_deck"):
            self.assertIn(document_type, entry_a)
            self.assertFalse(entry_a[document_type])

        gen_response = self.client.post(_generate_url(self.project_a.id))
        self.assertEqual(gen_response.status_code, 200)

        response = self.client.get(_status_url())
        results_by_project = {r["project"]: r for r in response.json()["results"]}

        self.assertTrue(results_by_project[self.project_a.id]["training_deck"])
        self.assertFalse(results_by_project[self.project_a.id]["system_test"])

        # project_b 完全不受影響。
        self.assertFalse(results_by_project[self.project_b.id]["training_deck"])


@override_settings(MEDIA_ROOT=_TEST_MEDIA_ROOT)
class PptxGenerateCsrfProtectionTests(TestCase):
    """產生端點是會改變狀態的寫入端點，training_deck／PPTX 分支必須跟
    DOCX 分支一樣受 CSRF 保護（絕不可 `@csrf_exempt`）。"""

    def setUp(self):
        self.user = User.objects.create_user(username="pm_pptx_csrf", password="pw12345")
        self.project = Project.objects.create(
            owner=self.user, customer_name="客戶", project_name="專案"
        )
        self.node = FeatureNode.objects.create(name="登入")
        FeatureNodeContent.objects.create(
            node=self.node, document_type="training_deck", bullets=["內容"]
        )
        ProjectFeatureSelection.objects.create(project=self.project, node=self.node)
        self.client = Client(enforce_csrf_checks=True)
        self.client.login(username="pm_pptx_csrf", password="pw12345")

    def _csrf_token(self):
        return self.client.get("/api/auth/csrf/").json()["csrfToken"]

    def test_generate_without_csrf_token_is_rejected(self):
        response = self.client.post(_generate_url(self.project.id))
        self.assertEqual(response.status_code, 403)
        self.assertEqual(GeneratedDocument.objects.count(), 0)

    def test_generate_with_primed_csrf_token_succeeds(self):
        token = self._csrf_token()
        response = self.client.post(
            _generate_url(self.project.id),
            HTTP_X_CSRFTOKEN=token,
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(GeneratedDocument.objects.count(), 1)
